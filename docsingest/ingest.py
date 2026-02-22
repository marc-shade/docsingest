import sys
import os
import re
import logging
from typing import Dict, Optional, Tuple, List
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords

# Download necessary NLTK resources
try:
    nltk.download('punkt', quiet=True)
    nltk.download('stopwords', quiet=True)
except Exception as e:
    print(f"Warning: Could not download NLTK resources: {e}", file=sys.stderr)

# Add the parent directory to Python path
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

import logging
from typing import Dict, Optional, Tuple

# Try importing PIIDetector
try:
    from docsingest.pii_detector import PIIDetector
except ImportError:
    print("Warning: PIIDetector could not be imported", file=sys.stderr)
    PIIDetector = None

# Try importing compliance modules
try:
    from docsingest.compliance.cui_detector import CUIDetector
    from docsingest.compliance.enhanced_pii import EnhancedPIIDetector
    from docsingest.compliance.sanitizer import DocumentSanitizer
    from docsingest.compliance.export_control import ExportControlScreener
    from docsingest.compliance.audit_trail import AuditTrail, AuditEventType, AuditSeverity
    COMPLIANCE_AVAILABLE = True
except ImportError:
    COMPLIANCE_AVAILABLE = False

logger = logging.getLogger(__name__)

# Default compliance-focused prompt
DEFAULT_COMPLIANCE_PROMPT = """
You are a compliance and risk management AI assistant.
Your task is to:
1. Analyze documents for sensitive information
2. Identify potential compliance risks
3. Provide actionable recommendations
4. Ensure data privacy and protection standards are met
"""


def ingest(
    input_directory: str,
    agent_prompt: str = None,
    output_file: Optional[str] = None,
    pii_analysis: bool = True,
    verbose: bool = False,
    compress_content: bool = False,
    compression_level: float = 0.5,
    cui_scan: bool = False,
    sanitize: bool = False,
    export_control: bool = False,
    enhanced_pii: bool = False,
    audit_log_path: Optional[str] = None,
    compliance_report_path: Optional[str] = None,
) -> Tuple[str, str, Dict, Dict]:
    """
    Recursively ingest documents from the specified directory,
    capturing full text content while protecting PII.

    Args:
        input_directory: Path to the directory to ingest documents from
        agent_prompt: Optional agent prompt (for backward compatibility)
        output_file: Optional path to save the output
        pii_analysis: Enable/disable PII detection
        verbose: Enable detailed logging
        compress_content: Whether to compress document content
        compression_level: Level of content compression (0.0 to 1.0)
        cui_scan: Enable CUI detection per 32 CFR Part 2002
        sanitize: Enable document sanitization analysis
        export_control: Enable ITAR/EAR export control screening
        enhanced_pii: Enable defense-grade enhanced PII/PHI detection
        audit_log_path: Path for audit trail output
        compliance_report_path: Path for separate compliance report

    Returns:
        Tuple of (summary, document_tree, document_contents, pii_reports)
    """
    try:
        # Use default compliance prompt if not specified
        prompt = agent_prompt or DEFAULT_COMPLIANCE_PROMPT

        # Initialize audit trail if requested
        audit_trail = None
        if audit_log_path and COMPLIANCE_AVAILABLE:
            audit_trail = AuditTrail(log_path=audit_log_path)
            logger.info("Audit trail initialized: %s", audit_log_path)

        # Perform document ingestion
        document_context = ingest_documents(
            input_directory,
            verbose=verbose,
            compress_content=compress_content,
            compression_level=compression_level
        )

        # Add output file to context for writing
        document_context['output_file'] = output_file or 'document_context.md'

        # Perform PII analysis if enabled
        pii_reports = {}
        if pii_analysis:
            from .pii_detector import PIIDetector
            pii_detector = PIIDetector()

            # Analyze PII for each document
            for directory, docs in document_context.get('document_groups', {}).items():
                for doc in docs:
                    content = doc.get('sanitized_content', '')
                    pii_report = pii_detector.detect(content)
                    if pii_report.get('pii_detected'):
                        pii_reports[doc.get('relative_path', 'Unknown')] = pii_report

        # Run compliance modules if available and requested
        compliance_results = {}
        any_compliance = cui_scan or sanitize or export_control or enhanced_pii

        if any_compliance and COMPLIANCE_AVAILABLE:
            compliance_results = _run_compliance_pipeline(
                input_directory=input_directory,
                document_context=document_context,
                cui_scan=cui_scan,
                sanitize=sanitize,
                export_control=export_control,
                enhanced_pii=enhanced_pii,
                audit_trail=audit_trail,
            )
            document_context['compliance_results'] = compliance_results
        elif any_compliance and not COMPLIANCE_AVAILABLE:
            msg = (
                "Compliance features requested but compliance modules not available. "
                "Install with: pip install docsingest[compliance]"
            )
            # In defense mode, missing compliance modules is a hard error
            if cui_scan and sanitize and export_control and enhanced_pii:
                raise RuntimeError(
                    f"DEFENSE MODE ERROR: {msg} "
                    "Defense mode requires all compliance modules to be installed."
                )
            logger.warning(msg)

        # Write document context to file
        write_document_context(document_context)

        # Write compliance report if requested
        if compliance_report_path and compliance_results:
            _write_compliance_report(compliance_report_path, compliance_results)

        # Log audit events for document ingestion
        if audit_trail:
            for directory, docs in document_context.get('document_groups', {}).items():
                for doc in docs:
                    file_path = os.path.join(input_directory, doc.get('relative_path', ''))
                    audit_trail.log_document_ingest(
                        document_path=file_path,
                        file_type=doc.get('extension', 'unknown'),
                        tokens=doc.get('tokens', 0),
                    )

            # Log export event
            audit_trail.log_document_export(
                document_path=document_context.get('output_file', 'unknown'),
                export_format="markdown",
                destination=document_context.get('output_file', 'unknown'),
            )

        # Generate document tree
        document_tree = generate_document_tree(input_directory)

        # Build summary with compliance info
        summary = document_context.get('summary', 'No summary available')
        if compliance_results:
            summary += _build_compliance_summary(compliance_results)

        # Return comprehensive results
        return (
            summary,
            document_tree,
            document_context.get('document_groups', {}),
            pii_reports
        )

    except Exception as e:
        print(f"Error during document analysis: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc(file=sys.stderr)
        raise


def _run_compliance_pipeline(
    input_directory: str,
    document_context: Dict,
    cui_scan: bool,
    sanitize: bool,
    export_control: bool,
    enhanced_pii: bool,
    audit_trail: Optional[object] = None,
) -> Dict:
    """
    Run all enabled compliance modules across ingested documents.

    Returns dictionary with compliance results keyed by module name.
    """
    results = {
        'cui': {},
        'enhanced_pii': {},
        'sanitization': {},
        'export_control': {},
    }

    # Initialize detectors
    cui_detector = CUIDetector() if cui_scan else None
    pii_detector = EnhancedPIIDetector() if enhanced_pii else None
    sanitizer = DocumentSanitizer() if sanitize else None
    export_screener = ExportControlScreener() if export_control else None

    for directory, docs in document_context.get('document_groups', {}).items():
        for doc in docs:
            relative_path = doc.get('relative_path', 'Unknown')
            content = doc.get('sanitized_content', '')
            full_path = os.path.join(input_directory, relative_path)

            if not content:
                continue

            # CUI Detection
            if cui_detector:
                try:
                    cui_report = cui_detector.detect(content, filename=relative_path)
                    if cui_report.cui_detected or cui_report.classification_detected or cui_report.risk_score > 0:
                        results['cui'][relative_path] = {
                            'cui_detected': cui_report.cui_detected,
                            'classification_detected': cui_report.classification_detected,
                            'risk_score': cui_report.risk_score,
                            'findings_count': len(cui_report.findings),
                            'marking_deficiencies': cui_report.marking_deficiencies,
                            'handling_recommendations': cui_report.handling_recommendations,
                            'nist_controls': cui_report.nist_800_171_controls,
                            'summary': cui_report.summary,
                        }
                    if audit_trail:
                        audit_trail.log_compliance_scan(
                            document_path=full_path,
                            scan_type="cui",
                            findings_count=len(cui_report.findings),
                            risk_score=cui_report.risk_score,
                        )
                except Exception as e:
                    logger.error("CUI detection failed for %s: %s", relative_path, e)

            # Enhanced PII Detection
            if pii_detector:
                try:
                    pii_report = pii_detector.detect(content, filename=relative_path)
                    if pii_report.pii_detected:
                        results['enhanced_pii'][relative_path] = {
                            'pii_detected': True,
                            'total_findings': pii_report.total_findings,
                            'risk_score': pii_report.risk_score,
                            'detections_by_category': pii_report.detections_by_category,
                            'detections_by_regulation': pii_report.detections_by_regulation,
                            'remediation_actions': pii_report.remediation_actions,
                            'nist_controls': pii_report.nist_controls_applicable,
                            'summary': pii_report.summary,
                        }
                    if audit_trail:
                        audit_trail.log_compliance_scan(
                            document_path=full_path,
                            scan_type="pii",
                            findings_count=pii_report.total_findings,
                            risk_score=pii_report.risk_score,
                        )
                except Exception as e:
                    logger.error("Enhanced PII detection failed for %s: %s", relative_path, e)

            # Document Sanitization (file-based analysis)
            if sanitizer:
                try:
                    if os.path.exists(full_path):
                        san_report = sanitizer.analyze(full_path)
                    else:
                        san_report = sanitizer.analyze_text(content, filename=relative_path)
                    if san_report.findings:
                        results['sanitization'][relative_path] = {
                            'findings_count': len(san_report.findings),
                            'risk_score': san_report.risk_score,
                            'findings_by_type': san_report.findings_by_type,
                            'findings_by_severity': san_report.findings_by_severity,
                            'macros_detected': san_report.macros_detected,
                            'hidden_content_detected': san_report.hidden_content_detected,
                            'embedded_files': san_report.embedded_files_found,
                            'metadata_stripped': san_report.metadata_stripped,
                            'sha256': san_report.sha256_before,
                            'summary': san_report.summary,
                        }
                    if audit_trail:
                        audit_trail.log_compliance_scan(
                            document_path=full_path,
                            scan_type="sanitization",
                            findings_count=len(san_report.findings),
                            risk_score=san_report.risk_score,
                        )
                except Exception as e:
                    logger.error("Sanitization analysis failed for %s: %s", relative_path, e)

            # Export Control Screening
            if export_screener:
                try:
                    ec_report = export_screener.screen(content, filename=relative_path)
                    if ec_report.export_controlled:
                        results['export_control'][relative_path] = {
                            'export_controlled': True,
                            'itar_findings': ec_report.itar_findings,
                            'ear_findings': ec_report.ear_findings,
                            'risk_score': ec_report.risk_score,
                            'usml_categories': ec_report.usml_categories_referenced,
                            'eccn_patterns': ec_report.eccn_patterns_found,
                            'controlled_technologies': ec_report.controlled_technologies,
                            'classification_recommendations': ec_report.classification_recommendations,
                            'nist_controls': ec_report.nist_controls_applicable,
                            'summary': ec_report.summary,
                        }
                    if audit_trail:
                        audit_trail.log_compliance_scan(
                            document_path=full_path,
                            scan_type="export_control",
                            findings_count=len(ec_report.findings),
                            risk_score=ec_report.risk_score,
                        )
                except Exception as e:
                    logger.error("Export control screening failed for %s: %s", relative_path, e)

    return results


def _build_compliance_summary(compliance_results: Dict) -> str:
    """Build a compliance summary section for the output."""
    lines = ["\n\n## Defense Compliance Summary\n"]

    # CUI Summary
    cui_results = compliance_results.get('cui', {})
    if cui_results:
        cui_count = sum(1 for r in cui_results.values() if r.get('cui_detected'))
        max_risk = max((r.get('risk_score', 0) for r in cui_results.values()), default=0)
        lines.append(f"### CUI Detection")
        lines.append(f"- Documents with CUI markings: {cui_count}/{len(cui_results)}")
        lines.append(f"- Maximum risk score: {max_risk}/100")
        lines.append("")

    # Enhanced PII Summary
    pii_results = compliance_results.get('enhanced_pii', {})
    if pii_results:
        total_findings = sum(r.get('total_findings', 0) for r in pii_results.values())
        max_risk = max((r.get('risk_score', 0) for r in pii_results.values()), default=0)
        lines.append(f"### Enhanced PII/PHI Detection")
        lines.append(f"- Documents with PII: {len(pii_results)}")
        lines.append(f"- Total PII findings: {total_findings}")
        lines.append(f"- Maximum risk score: {max_risk}/100")
        lines.append("")

    # Sanitization Summary
    san_results = compliance_results.get('sanitization', {})
    if san_results:
        macros = sum(1 for r in san_results.values() if r.get('macros_detected'))
        hidden = sum(1 for r in san_results.values() if r.get('hidden_content_detected'))
        lines.append(f"### Document Sanitization")
        lines.append(f"- Documents with findings: {len(san_results)}")
        lines.append(f"- Documents with macros: {macros}")
        lines.append(f"- Documents with hidden content: {hidden}")
        lines.append("")

    # Export Control Summary
    ec_results = compliance_results.get('export_control', {})
    if ec_results:
        itar = sum(1 for r in ec_results.values() if r.get('itar_findings', 0) > 0)
        ear = sum(1 for r in ec_results.values() if r.get('ear_findings', 0) > 0)
        lines.append(f"### Export Control Screening")
        lines.append(f"- Documents with export controls: {len(ec_results)}")
        lines.append(f"- ITAR findings: {itar} documents")
        lines.append(f"- EAR findings: {ear} documents")
        lines.append("")

    if len(lines) == 1:
        return "\n\n## Defense Compliance Summary\nNo compliance findings detected.\n"

    return '\n'.join(lines)


def _write_compliance_report(report_path: str, compliance_results: Dict) -> None:
    """Write a detailed compliance report to a separate file."""
    try:
        with open(report_path, 'w', encoding='utf-8') as f:
            f.write("# Defense Compliance Report\n\n")
            f.write(f"Generated by docsingest v0.2.0\n\n")
            f.write("---\n\n")

            # CUI Findings
            cui_results = compliance_results.get('cui', {})
            if cui_results:
                f.write("## CUI Detection Results\n\n")
                for doc_path, result in sorted(cui_results.items()):
                    f.write(f"### {doc_path}\n")
                    f.write(f"- **CUI Detected**: {result.get('cui_detected', False)}\n")
                    f.write(f"- **Classification Detected**: {result.get('classification_detected', False)}\n")
                    f.write(f"- **Risk Score**: {result.get('risk_score', 0)}/100\n")
                    f.write(f"- **Findings**: {result.get('findings_count', 0)}\n")
                    if result.get('marking_deficiencies'):
                        f.write("- **Marking Deficiencies**:\n")
                        for deficiency in result['marking_deficiencies']:
                            f.write(f"  - {deficiency}\n")
                    if result.get('handling_recommendations'):
                        f.write("- **Handling Recommendations**:\n")
                        for rec in result['handling_recommendations']:
                            f.write(f"  - {rec}\n")
                    if result.get('nist_controls'):
                        f.write(f"- **NIST 800-171 Controls**: {', '.join(result['nist_controls'])}\n")
                    f.write("\n")

            # Enhanced PII Findings
            pii_results = compliance_results.get('enhanced_pii', {})
            if pii_results:
                f.write("## Enhanced PII/PHI Detection Results\n\n")
                for doc_path, result in sorted(pii_results.items()):
                    f.write(f"### {doc_path}\n")
                    f.write(f"- **Total Findings**: {result.get('total_findings', 0)}\n")
                    f.write(f"- **Risk Score**: {result.get('risk_score', 0)}/100\n")
                    if result.get('detections_by_category'):
                        f.write("- **By Category**:\n")
                        for cat, count in sorted(result['detections_by_category'].items()):
                            f.write(f"  - {cat}: {count}\n")
                    if result.get('detections_by_regulation'):
                        f.write("- **By Regulation**:\n")
                        for reg, count in sorted(result['detections_by_regulation'].items()):
                            f.write(f"  - {reg}: {count}\n")
                    if result.get('remediation_actions'):
                        f.write("- **Remediation Actions**:\n")
                        for action in result['remediation_actions']:
                            f.write(f"  - {action}\n")
                    f.write("\n")

            # Sanitization Findings
            san_results = compliance_results.get('sanitization', {})
            if san_results:
                f.write("## Document Sanitization Results\n\n")
                for doc_path, result in sorted(san_results.items()):
                    f.write(f"### {doc_path}\n")
                    f.write(f"- **Risk Score**: {result.get('risk_score', 0)}/100\n")
                    f.write(f"- **Macros Detected**: {result.get('macros_detected', False)}\n")
                    f.write(f"- **Hidden Content**: {result.get('hidden_content_detected', False)}\n")
                    f.write(f"- **SHA-256**: `{result.get('sha256', 'N/A')}`\n")
                    if result.get('findings_by_type'):
                        f.write("- **Findings by Type**:\n")
                        for ftype, count in sorted(result['findings_by_type'].items()):
                            f.write(f"  - {ftype}: {count}\n")
                    if result.get('metadata_stripped'):
                        f.write("- **Metadata Found**:\n")
                        for key, val in sorted(result['metadata_stripped'].items()):
                            f.write(f"  - {key}: {val}\n")
                    f.write("\n")

            # Export Control Findings
            ec_results = compliance_results.get('export_control', {})
            if ec_results:
                f.write("## Export Control Screening Results\n\n")
                for doc_path, result in sorted(ec_results.items()):
                    f.write(f"### {doc_path}\n")
                    f.write(f"- **ITAR Findings**: {result.get('itar_findings', 0)}\n")
                    f.write(f"- **EAR Findings**: {result.get('ear_findings', 0)}\n")
                    f.write(f"- **Risk Score**: {result.get('risk_score', 0)}/100\n")
                    if result.get('usml_categories'):
                        f.write(f"- **USML Categories**: {', '.join(result['usml_categories'])}\n")
                    if result.get('eccn_patterns'):
                        f.write(f"- **ECCN Patterns**: {', '.join(result['eccn_patterns'])}\n")
                    if result.get('controlled_technologies'):
                        f.write("- **Controlled Technologies**:\n")
                        for tech in result['controlled_technologies']:
                            f.write(f"  - {tech}\n")
                    if result.get('classification_recommendations'):
                        f.write("- **Classification Recommendations**:\n")
                        for rec in result['classification_recommendations']:
                            f.write(f"  - {rec}\n")
                    f.write("\n")

            f.write("---\n\n")
            f.write("*Report generated by docsingest defense compliance module.*\n")

        print(f"Compliance report written to {report_path}", file=sys.stderr)

    except Exception as e:
        print(f"Error writing compliance report: {e}", file=sys.stderr)


def ingest_documents(
    input_directory,
    verbose=False,
    compress_content=False,
    compression_level=0.5
):
    """
    Core document ingestion function

    Args:
        input_directory (str): Path to the directory to ingest documents from
        verbose (bool): Enable detailed logging
        compress_content (bool): Whether to compress document content
        compression_level (float): Level of content compression (0.0 to 1.0)

    Returns:
        dict: Comprehensive document analysis results
    """
    # Load ignore patterns
    ignore_patterns = load_ignore_patterns(input_directory)

    # Setup logging
    log_file = setup_logging() if verbose else None

    # Prepare document tracking
    document_groups = {}
    total_tokens = 0

    # Walk through directory
    for root, dirs, files in os.walk(input_directory):
        # Apply ignore patterns
        dirs[:] = [d for d in dirs if not any(
            re.search(pattern, os.path.join(root, d))
            for pattern in ignore_patterns
        )]
        files = [f for f in files if not any(
            re.search(pattern, os.path.join(root, f))
            for pattern in ignore_patterns
        )]

        # Process files
        for filename in files:
            file_path = os.path.join(root, filename)
            relative_path = os.path.relpath(file_path, input_directory)

            try:
                # Extract text based on file extension
                text_extractor = _get_text_extractor(filename)
                if text_extractor:
                    content = text_extractor(file_path)

                    # Optional compression
                    if compress_content:
                        content = semantic_compress_text(content, compression_level)

                    # Tokenize content
                    tokens = len(content.split())
                    total_tokens += tokens

                    # Group documents by directory
                    directory = os.path.dirname(relative_path) or 'Root'
                    if directory not in document_groups:
                        document_groups[directory] = []

                    document_groups[directory].append({
                        'filename': filename,
                        'relative_path': relative_path,
                        'extension': os.path.splitext(filename)[1],
                        'sanitized_content': content,
                        'tokens': tokens,
                        'compression_applied': compress_content
                    })

            except Exception as e:
                if verbose:
                    logging.error(f"Error processing {file_path}: {e}")

    # Prepare summary
    summary = f"""
## Document Ingestion Summary
- **Total Directories**: {len(document_groups)}
- **Total Files Processed**: {sum(len(docs) for docs in document_groups.values())}
- **Total Tokens**: {total_tokens}
- **Compression**: {'Enabled' if compress_content else 'Disabled'}
"""

    return {
        'summary': summary,
        'document_groups': document_groups,
        'total_tokens': total_tokens,
        'log_file': log_file
    }


def semantic_compress_text(text: str, compression_level: float = 0.5) -> str:
    """
    Compress text while preserving semantic meaning

    Args:
        text (str): Input text to compress
        compression_level (float): Compression ratio (0.0 to 1.0)

    Returns:
        str: Compressed text
    """
    if not text or compression_level >= 1.0:
        return text

    try:
        # Tokenize sentences
        sentences = sent_tokenize(text)

        # If fewer sentences than compression allows, return full text
        if len(sentences) <= 2:
            return text

        # Remove stop words and calculate sentence importance
        try:
            stop_words = set(stopwords.words('english'))
        except LookupError:
            # Fallback if stopwords are not available
            stop_words = set()

        def sentence_importance(sentence):
            # Tokenize and remove stop words
            words = word_tokenize(sentence.lower())
            meaningful_words = [word for word in words if word.isalnum() and word not in stop_words]

            # Score based on unique meaningful words and length
            unique_word_score = len(set(meaningful_words))
            length_score = len(meaningful_words)

            return unique_word_score * length_score

        # Score sentences
        sentence_scores = [(sent, sentence_importance(sent)) for sent in sentences]

        # Sort sentences by importance
        sorted_sentences = sorted(sentence_scores, key=lambda x: x[1], reverse=True)

        # Calculate number of sentences to keep
        num_sentences = max(2, int(len(sentences) * compression_level))

        # Select top sentences, preserving original order
        top_sentences = sorted(sorted_sentences[:num_sentences], key=lambda x: sentences.index(x[0]))
        compressed_text = ' '.join(sent for sent, _ in top_sentences)

        # Add ellipsis if text was compressed
        if len(top_sentences) < len(sentences):
            compressed_text += " ... [Compressed]"

        return compressed_text

    except Exception as e:
        print(f"Error in semantic compression: {e}", file=sys.stderr)
        return text


def extract_text_from_docx(file_path):
    """Extract text from .docx files"""
    try:
        import docx
        doc = docx.Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs if para.text])
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}", file=sys.stderr)
        return f"[Error extracting text from {os.path.basename(file_path)}]"

def extract_text_from_xlsx(file_path):
    """Extract text from .xlsx files"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True)
        text_content = []
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text_content.append(row_text)
        return '\n'.join(text_content)
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}", file=sys.stderr)
        return f"[Error extracting text from {os.path.basename(file_path)}]"

def extract_text_from_xls(file_path):
    """Extract text from .xls files"""
    try:
        import xlrd
        wb = xlrd.open_workbook(file_path)
        text_content = []
        for sheet in wb.sheets():
            for row in range(sheet.nrows):
                row_text = ' '.join(str(cell.value) for cell in sheet.row(row) if str(cell.value).strip())
                if row_text:
                    text_content.append(row_text)
        return '\n'.join(text_content)
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}", file=sys.stderr)
        return f"[Error extracting text from {os.path.basename(file_path)}]"

def extract_text_from_pdf(file_path):
    """Extract text from .pdf files"""
    try:
        import PyPDF2
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ''
            for page in reader.pages:
                text += page.extract_text() + '\n'
            return text
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}", file=sys.stderr)
        return f"[Error extracting text from {os.path.basename(file_path)}]"

def extract_text_from_pptx(file_path):
    """Extract text from .pptx files"""
    try:
        import pptx
        prs = pptx.Presentation(file_path)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        return '\n'.join(text_content)
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}", file=sys.stderr)
        return f"[Error extracting text from {os.path.basename(file_path)}]"

def extract_text_from_json(file_path):
    """Extract text from .json files"""
    try:
        import json
        with open(file_path, 'r') as file:
            data = json.load(file)
            return json.dumps(data, indent=2)
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}", file=sys.stderr)
        return f"[Error extracting text from {os.path.basename(file_path)}]"

def extract_text_from_csv(file_path):
    """Extract text from .csv files"""
    try:
        import csv
        with open(file_path, 'r') as file:
            reader = csv.reader(file)
            return '\n'.join([','.join(row) for row in reader])
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}", file=sys.stderr)
        return f"[Error extracting text from {os.path.basename(file_path)}]"

def extract_text_from_xml(file_path):
    """Extract text from .xml files"""
    try:
        import xml.etree.ElementTree as ET
        tree = ET.parse(file_path)
        root = tree.getroot()

        def extract_text(element):
            text = element.text or ''
            for child in element:
                text += ' ' + extract_text(child)
            return text.strip()

        return extract_text(root)
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}", file=sys.stderr)
        return f"[Error extracting text from {os.path.basename(file_path)}]"

def extract_text_from_markdown(file_path):
    """Extract text from .md files"""
    try:
        import markdown
        with open(file_path, 'r') as file:
            text = file.read()
            # Convert markdown to plain text
            return markdown.markdown(text)
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}", file=sys.stderr)
        return f"[Error extracting text from {os.path.basename(file_path)}]"


def _get_text_extractor(filename: str):
    """
    Get the appropriate text extractor based on file extension.

    Args:
        filename (str): Name of the file to extract text from

    Returns:
        Callable or None: Text extraction function or None if no extractor found
    """
    ext = os.path.splitext(filename)[1].lower()

    text_extractors = {
        '.docx': extract_text_from_docx,
        '.doc': extract_text_from_docx,  # Fallback for older Word docs
        '.pdf': extract_text_from_pdf,
        '.xlsx': extract_text_from_xlsx,
        '.xls': extract_text_from_xls,
        '.pptx': extract_text_from_pptx,
        '.ppt': extract_text_from_pptx,  # Fallback for older PowerPoint
        '.json': extract_text_from_json,
        '.csv': extract_text_from_csv,
        '.xml': extract_text_from_xml,
        '.md': extract_text_from_markdown,
        '.txt': lambda path: open(path, 'r', encoding='utf-8', errors='replace').read()
    }

    # Skip system files and hidden files
    if filename.startswith('.'):
        return None

    return text_extractors.get(ext)


def setup_logging():
    """
    Set up logging configuration

    Returns:
        str: Path to the log file
    """
    log_directory = os.path.join(os.getcwd(), 'logs')
    os.makedirs(log_directory, exist_ok=True)

    log_path = os.path.join(log_directory, 'docsingest.log')

    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
        handlers=[
            logging.FileHandler(log_path),
            logging.StreamHandler()  # This will also print logs to console
        ]
    )

    return log_path

def write_document_context(document_context):
    """
    Write document context to a markdown file with comprehensive details

    Args:
        document_context (dict): Comprehensive document analysis results
    """
    try:
        # Use a default output path if not specified
        output_path = document_context.get('output_file', 'document_context.md')

        with open(output_path, 'w', encoding='utf-8') as f:
            # Write overall summary
            f.write("# Document Ingestion Report\n\n")

            # Write summary from document groups
            summary = document_context.get('summary', '')
            f.write(summary + "\n\n")

            # Get document groups
            document_groups = document_context.get('document_groups', {})

            # Compute total files
            total_files = sum(len(docs) for docs in document_groups.values())

            # Add total files to summary if not already present
            if 'Total Files' not in summary:
                f.write(f"## Summary\n")
                f.write(f"- **Total Files**: {total_files}\n")
                f.write(f"- **Total Tokens**: {document_context.get('total_tokens', 'N/A')}\n\n")

            # Sort directories for consistent output
            for directory, dir_docs in sorted(document_groups.items()):
                f.write(f"### Directory: `{directory or 'Root'}`\n")
                f.write(f"**Total Files**: {len(dir_docs)}\n\n")

                for doc in sorted(dir_docs, key=lambda x: x['filename']):
                    f.write(f"#### {doc.get('filename', 'Unknown')}\n")
                    f.write(f"- **Path**: `{doc.get('relative_path', 'N/A')}`\n")
                    f.write(f"- **Type**: {doc.get('extension', 'N/A')}\n")
                    f.write(f"- **Tokens**: {doc.get('tokens', 'N/A')}\n")
                    f.write(f"- **Compression Applied**: {doc.get('compression_applied', False)}\n\n")

                    # Full Content Section
                    f.write("##### Full Content\n")
                    f.write("```\n")

                    # Retrieve the full, untruncated content
                    content = doc.get('sanitized_content', 'No content available')

                    # Replace any problematic characters
                    content = content.replace('\x00', '')  # Remove null bytes

                    f.write(content)
                    f.write("\n```\n\n")

                    # Optional Compressed View
                    if doc.get('compression_applied'):
                        f.write("##### Compressed Content View\n")
                        f.write("```\n")

                        # Use the semantic compression function to generate a compressed view
                        compressed_content = semantic_compress_text(content)
                        f.write(compressed_content)
                        f.write("\n```\n\n")

            # Write compliance results if present
            compliance_results = document_context.get('compliance_results', {})
            if compliance_results:
                f.write(_build_compliance_summary(compliance_results))

        print(f"Document context written to {output_path}", file=sys.stderr)

    except Exception as e:
        print(f"Error in document context writing: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc(file=sys.stderr)
        sys.exit(1)


def generate_document_tree(directory):
    """
    Generate a simple document tree representation

    Args:
        directory (str): Root directory to generate tree for

    Returns:
        str: Text representation of document tree
    """
    import os

    tree = []
    for root, _, files in os.walk(directory):
        level = root.replace(directory, '').count(os.sep)
        indent = ' ' * 4 * level
        tree.append(f"{indent}{os.path.basename(root)}/")
        subindent = ' ' * 4 * (level + 1)
        for file in files:
            tree.append(f"{subindent}{file}")

    return '\n'.join(tree)


def load_ignore_patterns(directory: str) -> List[str]:
    """
    Load ignore patterns from .docsingest_ignore file.

    Args:
        directory (str): Base directory to search for .docsingest_ignore

    Returns:
        List[str]: List of ignore patterns
    """
    ignore_file = os.path.join(directory, '.docsingest_ignore')
    ignore_patterns = []

    if os.path.exists(ignore_file):
        with open(ignore_file, 'r') as f:
            ignore_patterns = [
                line.strip()
                for line in f
                if line.strip() and not line.startswith('#')
            ]

    return ignore_patterns


if __name__ == "__main__":
    import sys
    import argparse

    # Create argument parser
    parser = argparse.ArgumentParser(description="Document Ingestion Tool")
    parser.add_argument("input_directory", help="Directory to ingest documents from")
    parser.add_argument("--compress", action="store_true",
                        help="Enable content compression")
    parser.add_argument("--compression-level", type=float, default=0.5,
                        help="Compression level (0.0 to 1.0, default: 0.5)")
    parser.add_argument("--verbose", action="store_true",
                        help="Enable verbose logging")

    # Parse arguments
    args = parser.parse_args()

    # Call ingest function
    try:
        result = ingest(
            args.input_directory,
            verbose=args.verbose,
            compress_content=args.compress,
            compression_level=args.compression_level
        )
        print("Ingest completed successfully", file=sys.stderr)
    except Exception as e:
        print(f"Ingest failed: {e}", file=sys.stderr)
        sys.exit(1)
